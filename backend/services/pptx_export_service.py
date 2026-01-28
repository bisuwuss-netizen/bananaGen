# backend/services/pptx_export_service.py
"""
PPTX 导出服务 - 从 SlidePage 数据结构生成可编辑的 PowerPoint 文件
"""

import logging
from pathlib import Path
from typing import List, Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .render_schemas import SlidePage, StyleConfig

logger = logging.getLogger(__name__)


class PPTXExportService:
    """PPTX 导出服务 - 生成可编辑的 PowerPoint 文件"""

    # 16:9 幻灯片尺寸
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # 默认颜色配置
    DEFAULT_COLORS = {
        "primary": "#1e3a8a",
        "secondary": "#64748b",
        "accent": "#3b82f6",
        "text": "#1e293b",
        "background": "#f8fafc",
    }

    def __init__(self, style_config: Optional[StyleConfig] = None):
        """
        初始化导出服务
        
        Args:
            style_config: 风格配置
        """
        self.style = style_config
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        
        # 设置核心属性
        self._set_core_properties()

    def _set_core_properties(self):
        """设置文档元数据"""
        from datetime import datetime, timezone
        try:
            core = self.prs.core_properties
            now = datetime.now(timezone.utc)
            core.author = "banana-slides"
            core.last_modified_by = "banana-slides"
            core.created = now
            core.modified = now
            core.title = "职教课件"
            core.subject = "由 banana-slides 自动生成"
        except Exception as e:
            logger.warning(f"设置文档属性失败: {e}")

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """将十六进制颜色转换为 RGB 元组"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _get_color(self, color_key: str) -> RGBColor:
        """获取颜色值"""
        if self.style and self.style.color.get(color_key):
            hex_val = self.style.color[color_key]
        else:
            hex_val = self.DEFAULT_COLORS.get(color_key, "#000000")
        
        r, g, b = self._hex_to_rgb(hex_val)
        return RGBColor(r, g, b)

    def create_from_slides(
        self,
        slides: List[SlidePage],
        output_path: str,
        image_paths: Dict[str, str] = None
    ) -> str:
        """
        从 SlidePage 列表创建 PPTX
        
        Args:
            slides: 幻灯片页面列表
            output_path: 输出文件路径
            image_paths: 图片路径映射 {slot_id: file_path}
        
        Returns:
            输出文件路径
        """
        image_paths = image_paths or {}

        logger.info(f"[PPTXExport] 开始导出: {len(slides)} 页")

        for slide_data in slides:
            # 选择母版布局
            layout = self._get_layout_by_type(slide_data.slide_type)
            slide = self.prs.slides.add_slide(layout)

            # 根据布局类型渲染内容
            if slide_data.slide_type == "title":
                self._render_title_slide(slide, slide_data)
            elif slide_data.layout_id == "operation_steps":
                self._render_steps_slide(slide, slide_data, image_paths)
            elif slide_data.layout_id == "concept_comparison":
                self._render_comparison_slide(slide, slide_data, image_paths)
            elif slide_data.layout_id in ["title_bullets_right_img", "center_visual", "split_vertical"]:
                self._render_image_slide(slide, slide_data, image_paths)
            elif slide_data.layout_id == "grid_4":
                self._render_grid_slide(slide, slide_data, image_paths)
            else:
                self._render_default_slide(slide, slide_data, image_paths)

            # 添加演讲者备注
            if slide_data.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data.speaker_notes

            # 添加教法环节标签
            if slide_data.pedagogy_phase:
                self._add_phase_tag(slide, slide_data.pedagogy_phase)

        # 保存
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        
        logger.info(f"[PPTXExport] 导出完成: {output_path}")
        return output_path

    def _get_layout_by_type(self, slide_type: str):
        """根据页面类型选择母版布局"""
        # 使用空白布局以获得更大的自由度
        # 布局索引：0=标题, 1=标题和内容, 5=仅标题, 6=空白
        layout_idx_map = {
            "title": 6,       # 使用空白布局，自定义渲染
            "section": 6,
            "toc": 6,
            "concept": 6,
            "steps": 6,
            "comparison": 6,
            "exercises": 6,
            "summary": 6,
            "thank_you": 6,
            "intro": 6,
        }
        idx = layout_idx_map.get(slide_type, 6)
        return self.prs.slide_layouts[idx]

    def _render_title_slide(self, slide, slide_data: SlidePage):
        """渲染封面页"""
        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self._get_color("primary")
        p.alignment = PP_ALIGN.CENTER

        # 添加副标题（如果有）
        bullets = self._extract_bullets(slide_data)
        if bullets:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.8),
                Inches(11.333), Inches(1)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = bullets[0]
            p.font.size = Pt(24)
            p.font.color.rgb = self._get_color("secondary")
            p.alignment = PP_ALIGN.CENTER

        # 添加底部装饰条
        self._add_decoration_bar(slide)

    def _render_default_slide(self, slide, slide_data: SlidePage, image_paths: Dict[str, str]):
        """渲染默认布局（标题 + 要点）"""
        # 添加标题
        self._add_title(slide, slide_data.title)

        # 添加要点列表
        bullets = self._extract_bullets(slide_data)
        self._add_bullet_points(slide, bullets, Inches(0.8), Inches(1.8), Inches(12), Inches(5))

        # 添加底部装饰条
        self._add_decoration_bar(slide)

    def _render_image_slide(self, slide, slide_data: SlidePage, image_paths: Dict[str, str]):
        """渲染带图片的布局"""
        # 添加标题
        self._add_title(slide, slide_data.title)

        # 左侧要点（占 55%）
        bullets = self._extract_bullets(slide_data)
        self._add_bullet_points(slide, bullets, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5))

        # 右侧图片区域（占 40%）
        img_left = Inches(7.8)
        img_top = Inches(1.8)
        img_width = Inches(5)
        img_height = Inches(5)

        # 查找对应的图片
        slot_id = f"p{slide_data.index}_slot0"
        if slot_id in image_paths:
            try:
                slide.shapes.add_picture(
                    image_paths[slot_id],
                    img_left, img_top, img_width, img_height
                )
            except Exception as e:
                logger.warning(f"添加图片失败: {e}")
                self._add_image_placeholder(slide, img_left, img_top, img_width, img_height)
        else:
            self._add_image_placeholder(slide, img_left, img_top, img_width, img_height)

        # 添加底部装饰条
        self._add_decoration_bar(slide)

    def _render_steps_slide(self, slide, slide_data: SlidePage, image_paths: Dict[str, str]):
        """渲染操作步骤布局"""
        # 添加标题
        self._add_title(slide, slide_data.title)

        # 左侧图片区域
        img_left = Inches(0.8)
        img_top = Inches(1.8)
        img_width = Inches(5.5)
        img_height = Inches(5)

        slot_id = f"p{slide_data.index}_slot0"
        if slot_id in image_paths:
            try:
                slide.shapes.add_picture(
                    image_paths[slot_id],
                    img_left, img_top, img_width, img_height
                )
            except Exception as e:
                logger.warning(f"添加图片失败: {e}")
                self._add_image_placeholder(slide, img_left, img_top, img_width, img_height)
        else:
            self._add_image_placeholder(slide, img_left, img_top, img_width, img_height)

        # 右侧步骤列表
        bullets = self._extract_bullets(slide_data)
        self._add_numbered_steps(slide, bullets, Inches(6.8), Inches(1.8), Inches(6), Inches(5))

        # 添加底部装饰条
        self._add_decoration_bar(slide)

    def _render_comparison_slide(self, slide, slide_data: SlidePage, image_paths: Dict[str, str]):
        """渲染对比布局"""
        # 添加标题
        self._add_title(slide, slide_data.title)

        bullets = self._extract_bullets(slide_data)

        # 左侧对比项
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.8),
            Inches(5.8), Inches(5)
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(39, 174, 96)  # 绿色
        left_box.fill.fore_color.brightness = 0.8
        left_box.line.color.rgb = RGBColor(39, 174, 96)

        # 左侧标签
        if bullets:
            left_label = slide.shapes.add_textbox(
                Inches(1.3), Inches(5.5),
                Inches(5), Inches(1)
            )
            tf = left_label.text_frame
            p = tf.paragraphs[0]
            p.text = bullets[0] if bullets else "选项 A"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(39, 174, 96)
            p.alignment = PP_ALIGN.CENTER

        # 右侧对比项
        right_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(1.8),
            Inches(5.8), Inches(5)
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(192, 57, 43)  # 红色
        right_box.fill.fore_color.brightness = 0.8
        right_box.line.color.rgb = RGBColor(192, 57, 43)

        # 右侧标签
        if len(bullets) > 1:
            right_label = slide.shapes.add_textbox(
                Inches(7.3), Inches(5.5),
                Inches(5), Inches(1)
            )
            tf = right_label.text_frame
            p = tf.paragraphs[0]
            p.text = bullets[1]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
            p.alignment = PP_ALIGN.CENTER

        # 添加底部装饰条
        self._add_decoration_bar(slide)

    def _render_grid_slide(self, slide, slide_data: SlidePage, image_paths: Dict[str, str]):
        """渲染四宫格布局"""
        # 添加标题
        self._add_title(slide, slide_data.title)

        bullets = self._extract_bullets(slide_data)
        
        # 四宫格位置
        positions = [
            (Inches(0.8), Inches(1.8)),   # 左上
            (Inches(6.8), Inches(1.8)),   # 右上
            (Inches(0.8), Inches(4.2)),   # 左下
            (Inches(6.8), Inches(4.2)),   # 右下
        ]
        
        cell_width = Inches(5.8)
        cell_height = Inches(2.2)

        for i, (left, top) in enumerate(positions):
            # 添加单元格背景
            cell_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, cell_width, cell_height
            )
            cell_box.fill.solid()
            cell_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
            cell_box.line.color.rgb = RGBColor(226, 232, 240)

            # 添加单元格标签
            if i < len(bullets):
                label_box = slide.shapes.add_textbox(
                    left + Inches(0.3), top + cell_height - Inches(0.8),
                    cell_width - Inches(0.6), Inches(0.6)
                )
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = bullets[i]
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self._get_color("primary")
                p.alignment = PP_ALIGN.CENTER

        # 添加底部装饰条
        self._add_decoration_bar(slide)

    def _add_title(self, slide, title: str):
        """添加标题"""
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5),
            Inches(12), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self._get_color("primary")

        # 添加标题下划线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.45),
            Inches(0.8), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._get_color("accent")
        line.line.fill.background()

    def _add_bullet_points(self, slide, items: List[str], left, top, width, height, max_chars_per_slide: int = 300):
        """添加列表项，带字号自适应"""
        if not items:
            return

        # 创建文本框
        body_box = slide.shapes.add_textbox(left, top, width, height)
        tf = body_box.text_frame
        tf.word_wrap = True

        # 计算总字数以决定字号
        total_chars = sum(len(item) for item in items)

        if total_chars > 400:
            font_size = Pt(14)
        elif total_chars > 300:
            font_size = Pt(16)
        elif total_chars > 200:
            font_size = Pt(18)
        else:
            font_size = Pt(22)

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 添加项目符号
            p.text = f"• {item}"
            p.font.size = font_size
            p.font.color.rgb = self._get_color("text")
            p.space_after = Pt(8)
            p.level = 0

    def _add_numbered_steps(self, slide, items: List[str], left, top, width, height):
        """添加编号步骤列表"""
        if not items:
            return

        # 创建文本框
        body_box = slide.shapes.add_textbox(left, top, width, height)
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"{i + 1}. {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self._get_color("text")
            p.space_after = Pt(12)

    def _add_image_placeholder(self, slide, left, top, width, height):
        """添加图片占位符"""
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(243, 244, 246)
        placeholder.line.color.rgb = RGBColor(209, 213, 219)
        placeholder.line.dash_style = 2  # 虚线

        # 添加占位文字
        textbox = slide.shapes.add_textbox(
            left + width / 4, top + height / 2 - Inches(0.2),
            width / 2, Inches(0.5)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[配图区域]"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(156, 163, 175)
        p.alignment = PP_ALIGN.CENTER

    def _add_decoration_bar(self, slide):
        """添加底部装饰条"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.4),
            self.SLIDE_WIDTH, Inches(0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._get_color("accent")
        bar.line.fill.background()

    def _add_phase_tag(self, slide, phase_name: str):
        """添加教法环节标签"""
        tag_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(11.5), Inches(0.3),
            Inches(1.6), Inches(0.4)
        )
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = self._get_color("primary")
        tag_box.fill.fore_color.brightness = 0.2
        tag_box.line.fill.background()

        # 添加标签文字
        text_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(0.35),
            Inches(1.6), Inches(0.35)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase_name
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def _extract_bullets(self, page: SlidePage) -> List[str]:
        """提取要点列表"""
        bullets = []
        for elem in page.elements:
            if elem.type == "bullets":
                if isinstance(elem.content, dict):
                    bullets.extend(elem.content.get("items", []))
                elif isinstance(elem.content, list):
                    bullets.extend(elem.content)
        return bullets
